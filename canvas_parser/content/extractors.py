import json
from pathlib import Path

from canvas_parser.content.legacy_office import extract_legacy_office_text
from canvas_parser.content.ocr import ocr_image_file


PARSEABLE_MIME_TYPES = {
    'application/pdf': 'pdf',
    'text/plain': 'text',
    'text/markdown': 'text',
    'text/csv': 'text',
    'application/json': 'json',
    'application/msword': 'doc',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'image/png': 'image',
    'image/jpeg': 'image',
    'image/jpg': 'image',
    'image/gif': 'image',
    'image/webp': 'image',
}


def detect_extractor(content_type='', filename=''):
    mime = str(content_type or '').casefold().split(';')[0].strip()
    if mime in PARSEABLE_MIME_TYPES:
        return PARSEABLE_MIME_TYPES[mime]

    suffix = Path(str(filename or '')).suffix.casefold()
    mapping = {
        '.pdf': 'pdf',
        '.txt': 'text',
        '.md': 'text',
        '.csv': 'text',
        '.json': 'json',
        '.doc': 'doc',
        '.docx': 'docx',
        '.ppt': 'ppt',
        '.pptx': 'pptx',
        '.xlsx': 'xlsx',
        '.ipynb': 'ipynb',
        '.tex': 'text',
        '.py': 'text',
        '.m': 'text',
        '.r': 'text',
        '.c': 'text',
        '.cpp': 'text',
        '.h': 'text',
        '.java': 'text',
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.gif': 'image',
        '.webp': 'image',
    }
    return mapping.get(suffix, '')


def sniff_extractor_from_path(path):
    """Infer extractor from downloaded bytes when Canvas omits extension/content-type."""
    file_path = Path(path)
    if not file_path.is_file():
        return ''
    try:
        head = file_path.read_bytes()[:8]
    except OSError:
        return ''
    if head.startswith(b'%PDF'):
        return 'pdf'
    if head.startswith(b'PK\x03\x04'):
        return 'docx'
    return ''


def resolve_extractor_kind(content_type='', filename='', path=None):
    kind = detect_extractor(content_type, filename)
    if kind:
        return kind
    if path is not None:
        kind = sniff_extractor_from_path(path)
        if kind:
            return kind
    lowered = str(filename or '').casefold()
    if lowered.endswith('.pdf'):
        return 'pdf'
    return ''


def extract_text_from_file(path, extractor_kind, build_pdf_pages=None, fileid=None):
    path = Path(path)
    if not path.exists():
        return {'text': '', 'pages': []}

    if extractor_kind == 'pdf' and build_pdf_pages:
        pages = build_pdf_pages(str(path), str(fileid or path.name))
        text = '\n\n'.join(page.get('text', '') for page in pages if page.get('text'))
        return {'text': text, 'pages': pages}

    if extractor_kind in {'text', 'json'}:
        text = path.read_text(encoding='utf-8', errors='replace')
        return {'text': text, 'pages': []}

    if extractor_kind in {'doc', 'ppt'}:
        text = extract_legacy_office_text(path, suffix=path.suffix)
        return {'text': text, 'pages': []}

    if extractor_kind == 'image':
        text = ocr_image_file(path)
        if not text:
            text = f"[Image file: {path.name}]"
        return {'text': text, 'pages': []}

    if extractor_kind == 'docx':
        try:
            from docx import Document
        except ModuleNotFoundError:
            return {'text': '', 'pages': []}
        document = Document(str(path))
        text = '\n'.join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
        return {'text': text, 'pages': []}

    if extractor_kind == 'pptx':
        try:
            from pptx import Presentation
        except ModuleNotFoundError:
            return {'text': '', 'pages': []}
        presentation = Presentation(str(path))
        chunks = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                text = getattr(shape, 'text', '')
                if text:
                    slide_text.append(text)
            if slide_text:
                chunks.append(f"Slide {slide_index}\n" + '\n'.join(slide_text))
        text = '\n\n'.join(chunks)
        return {'text': text, 'pages': []}

    if extractor_kind == 'xlsx':
        try:
            from openpyxl import load_workbook
        except ModuleNotFoundError:
            return {'text': '', 'pages': []}
        workbook = load_workbook(str(path), read_only=True, data_only=True)
        rows = []
        for sheet in workbook.worksheets:
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell not in (None, '')]
                if cells:
                    sheet_rows.append('\t'.join(cells))
            if sheet_rows:
                rows.append(f"Sheet {sheet.title}\n" + '\n'.join(sheet_rows))
        workbook.close()
        text = '\n\n'.join(rows)
        return {'text': text, 'pages': []}

    if extractor_kind == 'ipynb':
        try:
            notebook = json.loads(path.read_text(encoding='utf-8', errors='replace'))
        except (json.JSONDecodeError, OSError):
            return {'text': '', 'pages': []}
        chunks = []
        for cell in notebook.get('cells', []):
            source = cell.get('source', [])
            if isinstance(source, list):
                source = ''.join(source)
            if source:
                chunks.append(str(source))
        return {'text': '\n\n'.join(chunks), 'pages': []}

    return {'text': '', 'pages': []}
